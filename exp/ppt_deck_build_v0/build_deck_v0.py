# -*- coding: utf-8 -*-
from __future__ import unicode_literals

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
ALIGN = ROOT / "align"
OUT_DIR = ROOT / "generated_pptx_test"
TEMPLATE = Path(r"D:\毕设\latex-for-zju-master\latex-for-zju-master\zjuslides.pptx")
FINAL_FIGS = Path(
    r"D:\毕设\latex-for-zju-master\latex-for-zju-master\figure\thesis_figures"
)
PROJECT_FIGS = ROOT / "gaokaollm_bench" / "outputs" / "thesis_figures"
GENERATED_ASSETS = Path(__file__).resolve().parent / "generated_assets"
DECK_PATH = OUT_DIR / "gaokaollm_defense_deck_v0.pptx"
MANIFEST_PATH = ALIGN / "ppt_deck_build_manifest_v0.md"
USE_VAI01_ON_MAIN = True
USE_ARCHITECTURE_IMAGE_ON_MAIN = False
SHOW_SOURCE_LINES = False

SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

COLORS = {
    "zju_blue": RGBColor(28, 64, 125),
    "light_blue": RGBColor(218, 235, 247),
    "pale_blue": RGBColor(238, 246, 252),
    "zju_red": RGBColor(146, 28, 36),
    "ink": RGBColor(35, 45, 55),
    "muted": RGBColor(92, 104, 115),
    "line": RGBColor(205, 213, 222),
    "pale": RGBColor(247, 249, 252),
    "blue": RGBColor(0, 114, 178),
    "orange": RGBColor(230, 159, 0),
    "green": RGBColor(0, 158, 115),
    "white": RGBColor(255, 255, 255),
}

FONT_CJK = "Microsoft YaHei"
FONT_SERIF = "SimSun"


def remove_all_slides(prs):
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def add_slide(prs):
    layout = (
        prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    )
    return prs.slides.add_slide(layout)


def add_box(
    slide,
    x,
    y,
    w,
    h,
    text,
    fill=None,
    line=None,
    font_size=16,
    bold=False,
    font_color=None,
    align=PP_ALIGN.CENTER,
    radius=True,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or COLORS["white"]
    shape.line.color.rgb = line or COLORS["line"]
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_CJK
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = font_color or COLORS["ink"]
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=15,
    bold=False,
    color=None,
    align=PP_ALIGN.LEFT,
    font=FONT_CJK,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]
    return box


def add_bullets(slide, x, y, w, h, items, size=14, color=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT_CJK
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["ink"]
        p.space_after = Pt(4)
    return box


def add_page_chrome(slide, watermark=True):
    add_box(
        slide,
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(7.5),
        "",
        fill=COLORS["white"],
        line=COLORS["white"],
        radius=False,
    )
    add_box(
        slide,
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(0.68),
        "",
        fill=COLORS["light_blue"],
        line=COLORS["light_blue"],
        radius=False,
    )
    add_box(
        slide,
        Inches(0),
        Inches(7.08),
        Inches(10),
        Inches(0.42),
        "",
        fill=COLORS["zju_blue"],
        line=COLORS["zju_blue"],
        radius=False,
    )
    if watermark:
        add_text(
            slide,
            Inches(6.25),
            Inches(1.05),
            Inches(3.1),
            Inches(1.0),
            "ZJU",
            size=48,
            bold=True,
            color=COLORS["pale_blue"],
            align=PP_ALIGN.RIGHT,
        )
        add_text(
            slide,
            Inches(7.0),
            Inches(6.92),
            Inches(2.7),
            Inches(0.18),
            "Zhejiang University",
            size=7.5,
            color=COLORS["white"],
            align=PP_ALIGN.RIGHT,
        )


def add_title(slide, title, subtitle=None, num=None):
    add_page_chrome(slide)
    title_size = (
        15
        if len(title) > 30
        else 16
        if len(title) > 24
        else 18
        if len(title) > 18
        else 20
    )
    add_text(
        slide,
        Inches(0.55),
        Inches(0.17),
        Inches(8.7),
        Inches(0.38),
        title,
        size=title_size,
        bold=True,
        color=COLORS["zju_blue"],
    )
    if subtitle:
        add_text(
            slide,
            Inches(0.58),
            Inches(0.72),
            Inches(8.75),
            Inches(0.24),
            subtitle,
            size=10,
            color=COLORS["muted"],
        )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.63), Inches(8.9), Pt(1.0)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["zju_blue"]
    line.line.color.rgb = COLORS["zju_blue"]


def add_takeaway(slide, text, y=Inches(6.32), color=None):
    return add_box(
        slide,
        Inches(0.75),
        y,
        Inches(8.5),
        Inches(0.42),
        text,
        fill=RGBColor(255, 251, 235),
        line=color or COLORS["orange"],
        font_size=12.5,
        bold=True,
    )


def add_source(slide, source):
    if not SHOW_SOURCE_LINES:
        return
    if source:
        add_text(
            slide,
            Inches(0.48),
            Inches(6.98),
            Inches(8.7),
            Inches(0.22),
            "Source: " + source,
            size=7.5,
            color=COLORS["muted"],
        )


def add_arrow(slide, x, y, w=Inches(0.34), h=Inches(0.22), color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color or COLORS["line"]
    shape.line.color.rgb = color or COLORS["line"]
    return shape


def image_path(name):
    candidates = []
    for base in (FINAL_FIGS, PROJECT_FIGS):
        for ext in (".png", ".jpg", ".jpeg"):
            candidates.append(base / (name + ext))
    for path in candidates:
        if path.exists():
            return path
    return None


def generated_asset(name):
    path = GENERATED_ASSETS / name
    return path if path.exists() else None


def cropped_figure(name):
    return generated_asset("crop_" + name + ".png") or image_path(name)


def vai01_asset():
    return generated_asset("v_ai01_openrouter_icu.png") or generated_asset(
        "v_ai01_algorithm_flow.png"
    )


def vai01_asset_label():
    if generated_asset("v_ai01_openrouter_icu.png"):
        return "v_ai01_openrouter_icu.png"
    if generated_asset("v_ai01_algorithm_flow.png"):
        return "v_ai01_algorithm_flow.png"
    return ""


def add_image_fit(slide, img_path, x, y, w, h):
    if not img_path or not Path(img_path).exists():
        add_box(
            slide,
            x,
            y,
            w,
            h,
            "素材待补\n" + (str(img_path) if img_path else ""),
            fill=RGBColor(250, 250, 250),
            line=COLORS["orange"],
            font_size=12,
        )
        return False
    with Image.open(str(img_path)) as im:
        iw, ih = im.size
    img_ratio = float(iw) / float(ih)
    box_ratio = float(w) / float(h)
    if img_ratio >= box_ratio:
        pic_w = w
        pic_h = int(w / img_ratio)
    else:
        pic_h = h
        pic_w = int(h * img_ratio)
    pic_x = x + int((w - pic_w) / 2)
    pic_y = y + int((h - pic_h) / 2)
    slide.shapes.add_picture(str(img_path), pic_x, pic_y, width=pic_w, height=pic_h)
    return True


def add_table_like(slide, x, y, w, row_h, headers, rows, col_ratios=None, font_size=11):
    col_ratios = col_ratios or [1.0 / len(headers)] * len(headers)
    col_ws = [int(w * r) for r in col_ratios]
    for idx, header in enumerate(headers):
        cx = x + sum(col_ws[:idx])
        add_box(
            slide,
            cx,
            y,
            col_ws[idx],
            row_h,
            header,
            fill=COLORS["zju_blue"],
            line=COLORS["zju_blue"],
            font_size=font_size,
            bold=True,
            font_color=COLORS["white"],
            radius=False,
        )
    for r_idx, row in enumerate(rows):
        for c_idx, cell in enumerate(row):
            cx = x + sum(col_ws[:c_idx])
            cy = y + row_h * (r_idx + 1)
            fill = RGBColor(248, 250, 252) if r_idx % 2 == 0 else COLORS["white"]
            add_box(
                slide,
                cx,
                cy,
                col_ws[c_idx],
                row_h,
                cell,
                fill=fill,
                line=COLORS["line"],
                font_size=font_size,
                radius=False,
            )


def add_pipeline(
    slide,
    labels,
    y,
    x0=Inches(0.58),
    box_w=Inches(1.25),
    box_h=Inches(0.58),
    gap=Inches(0.18),
    colors=None,
    font_size=11,
):
    colors = colors or [COLORS["blue"]] * len(labels)
    for i, label in enumerate(labels):
        x = x0 + i * (box_w + gap)
        add_box(
            slide,
            x,
            y,
            box_w,
            box_h,
            label,
            fill=COLORS["white"],
            line=colors[i],
            font_size=font_size,
            bold=True,
        )
        if i < len(labels) - 1:
            add_arrow(
                slide,
                x + box_w + Inches(0.02),
                y + Inches(0.19),
                w=Inches(0.14),
                h=Inches(0.18),
                color=COLORS["line"],
            )


def make_cover(prs):
    slide = add_slide(prs)
    add_box(
        slide,
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(7.5),
        "",
        fill=COLORS["white"],
        line=COLORS["white"],
        radius=False,
    )
    add_box(
        slide,
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(0.18),
        "",
        fill=COLORS["zju_blue"],
        line=COLORS["zju_blue"],
        radius=False,
    )
    add_text(
        slide,
        Inches(0.85),
        Inches(1.10),
        Inches(8.3),
        Inches(0.55),
        "本科毕业设计答辩",
        size=22,
        bold=True,
        color=COLORS["zju_red"],
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        Inches(0.95),
        Inches(2.08),
        Inches(8.1),
        Inches(0.9),
        "大模型驱动的高考志愿推荐系统设计与实现",
        size=29,
        bold=True,
        color=COLORS["zju_blue"],
        align=PP_ALIGN.CENTER,
    )
    add_box(
        slide,
        Inches(1.75),
        Inches(3.35),
        Inches(6.5),
        Inches(1.48),
        "学生：潘臻琦    学号：3210102495\n指导教师：胡天磊\n学院：计算机学院    专业：计算机科学与技术\n2026年5月18日",
        fill=RGBColor(248, 250, 252),
        line=COLORS["line"],
        font_size=15,
    )
    add_text(
        slide,
        Inches(1.1),
        Inches(6.72),
        Inches(7.8),
        Inches(0.25),
        "答辩主线：事实可信 -> 偏好澄清 -> 可复查验证",
        size=11,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
    )
    return slide


def add_agenda_slide(prs):
    slide = add_slide(prs)
    add_page_chrome(slide)
    add_text(
        slide,
        Inches(0.9),
        Inches(0.95),
        Inches(8.2),
        Inches(0.45),
        "目录",
        size=28,
        bold=True,
        color=COLORS["zju_blue"],
        align=PP_ALIGN.CENTER,
    )
    items = [
        ("1", "研究背景与问题定义"),
        ("2", "系统与方法设计"),
        ("3", "实验验证与性能分析"),
        ("4", "总结与展望"),
    ]
    for idx, (num, text) in enumerate(items):
        y = Inches(2.0 + idx * 0.82)
        add_box(
            slide,
            Inches(2.0),
            y,
            Inches(0.44),
            Inches(0.44),
            num,
            fill=COLORS["zju_blue"],
            line=COLORS["zju_blue"],
            font_size=15,
            bold=True,
            font_color=COLORS["white"],
            radius=False,
        )
        add_text(
            slide,
            Inches(2.75),
            y - Inches(0.02),
            Inches(5.4),
            Inches(0.5),
            text,
            size=18,
            bold=True,
            color=COLORS["ink"],
        )
    return slide


def add_section_divider(prs, num, title):
    slide = add_slide(prs)
    add_page_chrome(slide)
    add_text(
        slide,
        Inches(1.1),
        Inches(2.82),
        Inches(1.0),
        Inches(0.42),
        num,
        size=20,
        bold=True,
        color=COLORS["zju_blue"],
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        Inches(2.05),
        Inches(2.75),
        Inches(6.6),
        Inches(0.62),
        title,
        size=25,
        bold=True,
        color=COLORS["zju_blue"],
        align=PP_ALIGN.LEFT,
    )
    return slide


def slide_problem_combined(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "事实可召回，但真实底线需要通过取舍被澄清",
        "高考志愿推荐的难点不只是找信息，而是识别用户愿意放宽什么、坚持什么",
        "04",
    )
    add_box(
        slide,
        Inches(0.75),
        Inches(1.35),
        Inches(2.1),
        Inches(0.72),
        "显式话语\n分数 / 专业 / 地域",
        fill=RGBColor(241, 248, 255),
        line=COLORS["blue"],
        font_size=14,
        bold=True,
    )
    add_arrow(
        slide,
        Inches(2.95),
        Inches(1.58),
        w=Inches(0.35),
        h=Inches(0.18),
        color=COLORS["line"],
    )
    add_box(
        slide,
        Inches(3.35),
        Inches(1.35),
        Inches(2.1),
        Inches(0.72),
        "事实候选\n学校 / 专业 / 分数线",
        fill=COLORS["white"],
        line=COLORS["line"],
        font_size=14,
        bold=True,
    )
    add_arrow(
        slide,
        Inches(5.55),
        Inches(1.58),
        w=Inches(0.35),
        h=Inches(0.18),
        color=COLORS["line"],
    )
    add_box(
        slide,
        Inches(5.95),
        Inches(1.35),
        Inches(2.1),
        Inches(0.72),
        "隐藏底线\n预算 / 风险 / 距离",
        fill=RGBColor(255, 248, 240),
        line=COLORS["orange"],
        font_size=14,
        bold=True,
    )
    add_arrow(
        slide,
        Inches(8.15),
        Inches(1.58),
        w=Inches(0.35),
        h=Inches(0.18),
        color=COLORS["line"],
    )
    add_box(
        slide,
        Inches(8.45),
        Inches(1.35),
        Inches(0.8),
        Inches(0.72),
        "解释",
        fill=RGBColor(241, 250, 246),
        line=COLORS["green"],
        font_size=13,
        bold=True,
    )
    add_box(
        slide,
        Inches(1.15),
        Inches(2.95),
        Inches(7.7),
        Inches(1.2),
        "核心判断：信息入口能解决“查得到”，但不能自动回答\n用户是否愿意为更远地域、更高学费或更高风险放宽另一条约束。",
        fill=RGBColor(248, 250, 252),
        line=COLORS["line"],
        font_size=18,
        bold=True,
    )
    add_takeaway(slide, "本设计要把隐藏取舍变成可提问、可反馈、可回放的证据链。")
    add_source(slide, "fact_ledger_v0.md §2; PPT_storyboard_v0.md Slides 2-3")


def slide_fact_boundary_compact(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "事实边界由证据层控制，模型负责语义归一和表达组织",
        "学校、专业、分数、学费来自结构化数据；偏好取舍来自用户反馈",
        "08",
    )
    add_pipeline(
        slide,
        ["用户表达", "语义归一", "证据层检索", "事实候选", "解释输出"],
        Inches(1.72),
        x0=Inches(0.85),
        box_w=Inches(1.5),
        box_h=Inches(0.7),
        colors=[
            COLORS["line"],
            COLORS["green"],
            COLORS["blue"],
            COLORS["blue"],
            COLORS["zju_blue"],
        ],
        font_size=12,
    )
    add_box(
        slide,
        Inches(1.1),
        Inches(3.15),
        Inches(3.55),
        Inches(1.2),
        "LLM 不生成事实\n只做语义归一、探测规划和表达组织",
        fill=RGBColor(241, 250, 246),
        line=COLORS["green"],
        font_size=16,
        bold=True,
    )
    add_box(
        slide,
        Inches(5.3),
        Inches(3.15),
        Inches(3.55),
        Inches(1.2),
        "证据层给出候选\n学校 / 专业 / 分数线 / 位次 / 学费",
        fill=RGBColor(241, 248, 255),
        line=COLORS["blue"],
        font_size=16,
        bold=True,
    )
    add_takeaway(
        slide, "安全口径：事实由数据给出，取舍由用户确认，解释回到可查证证据。"
    )
    add_source(slide, "fact_ledger_v0.md §2, §4, §5, §8")


def slide_s12_state_combined(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "A/B 取舍把反馈写回偏好状态，而不是改写底层事实查询",
        "Pareto 候选负责呈现边际取舍，BT posterior 影响后续选轴、评分与解释",
        "13",
    )
    add_box(
        slide,
        Inches(0.7),
        Inches(1.25),
        Inches(2.15),
        Inches(1.35),
        "候选 A\n更稳妥\n专业匹配较弱\n学费更低",
        fill=RGBColor(241, 248, 255),
        line=COLORS["blue"],
        font_size=12.5,
        bold=True,
    )
    add_box(
        slide,
        Inches(3.05),
        Inches(1.25),
        Inches(2.15),
        Inches(1.35),
        "候选 B\n专业更匹配\n风险略高\n需放宽预算",
        fill=RGBColor(255, 248, 240),
        line=COLORS["orange"],
        font_size=12.5,
        bold=True,
    )
    add_box(
        slide,
        Inches(5.5),
        Inches(1.25),
        Inches(1.55),
        Inches(1.35),
        "反馈\nPrefer A\nPrefer B\nUncertain",
        fill=COLORS["white"],
        line=COLORS["line"],
        font_size=11.5,
    )
    add_box(
        slide,
        Inches(7.35),
        Inches(1.25),
        Inches(1.85),
        Inches(1.35),
        "后验状态\n地域 / 专业\n学费 / 风险",
        fill=RGBColor(241, 250, 246),
        line=COLORS["green"],
        font_size=12.5,
        bold=True,
    )
    add_image_fit(
        slide,
        cropped_figure("fig_5_2_runtime_state_machine"),
        Inches(0.85),
        Inches(3.02),
        Inches(8.3),
        Inches(2.85),
    )
    add_takeaway(slide, "反馈改变偏好状态和解释路径，不直接改写 SQL、query 或 filter。")
    add_source(slide, "fact_ledger_v0.md §5; fig_5_2_runtime_state_machine")


def slide_conclusion_combined(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "本设计初步支持“先问清底线，再给出可复查推荐”的技术路径",
        "当前结论是机制验证，不是生产级有效性或真实用户研究结论",
        "18",
    )
    add_box(
        slide,
        Inches(0.8),
        Inches(1.18),
        Inches(8.4),
        Inches(1.35),
        "工作总结\n1. 将高考志愿推荐组织为事实候选、主动澄清、在线后验和解释输出的闭环。\n2. 用受控画像与消融实验初步验证主动探测和后验追踪的机制价值。\n3. 所有结论回到事实边界与用户反馈，不把生成内容当作事实来源。",
        fill=COLORS["white"],
        line=COLORS["line"],
        font_size=13.5,
        bold=False,
        align=PP_ALIGN.LEFT,
    )
    add_box(
        slide,
        Inches(1.1),
        Inches(3.28),
        Inches(7.8),
        Inches(1.95),
        "研究边界与后续工作\n- 需要真实学生/专家评估，当前测试不替代真实用户研究。\n- 需要跨省份、跨年份和更多证据维度的泛化验证。\n- 需要进一步增强用户控制、隐私合规和解释透明度。",
        fill=RGBColor(248, 250, 252),
        line=COLORS["line"],
        font_size=14.5,
        bold=False,
        align=PP_ALIGN.LEFT,
    )
    add_takeaway(slide, "答辩口径：早期支持机制有效，后续仍需真实场景验证。")
    add_source(slide, "fact_ledger_v0.md §1, §8; 07-conclusion.tex:28-38")


def slide_s02(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "现有志愿填报方案能提供信息和建议，但隐藏底线仍需要被证据化澄清",
        "信息入口能解决查询问题，难点在于澄清真实取舍",
        "02",
    )
    headers = ["传统信息/筛选入口", "AI empowered 建议生成", "本文关注点"]
    rows = [
        ["查信息、筛条件、看规则", "问答式解释、生成建议", "证据化多轮澄清隐藏底线"],
        ["信息可得", "表达更自然", "事实边界 + A/B 取舍 + 后验状态"],
        ["依赖显式输入", "更会组织表达", "主动澄清真实取舍"],
    ]
    add_table_like(
        slide,
        Inches(0.7),
        Inches(1.35),
        Inches(8.6),
        Inches(0.72),
        headers,
        rows,
        col_ratios=[0.31, 0.31, 0.38],
        font_size=13,
    )
    add_box(
        slide,
        Inches(1.0),
        Inches(5.45),
        Inches(8.0),
        Inches(0.56),
        "信息入口存在 != 用户真实底线已经被澄清",
        fill=RGBColor(255, 251, 235),
        line=COLORS["orange"],
        font_size=16,
        bold=True,
    )
    add_source(
        slide,
        "storyboard Slide 2; product-specific screenshots deferred until URL/date capture",
    )


def slide_s03(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "高考志愿填报的难点不是回答问题，而是从不完整表达中发现真实底线",
        "显式话语不等于完整真实偏好",
        "03",
    )
    add_box(
        slide,
        Inches(0.9),
        Inches(1.35),
        Inches(8.2),
        Inches(0.72),
        "显式话语：想去大城市、计算机相关、不要太贵",
        fill=RGBColor(242, 247, 255),
        line=COLORS["blue"],
        font_size=16,
    )
    add_arrow(
        slide,
        Inches(4.75),
        Inches(2.25),
        w=Inches(0.5),
        h=Inches(0.28),
        color=COLORS["line"],
    )
    add_box(
        slide,
        Inches(1.2),
        Inches(2.75),
        Inches(7.6),
        Inches(0.7),
        "事实约束：分数 / 位次 / 选科 / 学费 / 地域 / 专业",
        fill=COLORS["white"],
        line=COLORS["line"],
        font_size=15,
    )
    add_box(
        slide,
        Inches(1.65),
        Inches(4.0),
        Inches(6.7),
        Inches(1.2),
        "隐藏底线：哪些条件绝不能破？\n哪些收益值得放宽另一条约束？",
        fill=RGBColor(255, 248, 240),
        line=COLORS["orange"],
        font_size=18,
        bold=True,
    )
    add_box(
        slide,
        Inches(2.2),
        Inches(5.65),
        Inches(5.6),
        Inches(0.58),
        "推荐解释必须回到真实证据和用户反馈",
        fill=RGBColor(241, 250, 246),
        line=COLORS["green"],
        font_size=15,
    )
    add_source(slide, "fact_ledger_v0.md §2; 01-introduction.tex")


def slide_s04(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "本设计把一次性推荐改造成事实约束下的多轮偏好澄清系统",
        "事实候选、取舍问题、偏好状态和解释输出形成闭环",
        "04",
    )
    add_image_fit(
        slide,
        cropped_figure("fig_5_1_mas_workflow"),
        Inches(0.55),
        Inches(1.05),
        Inches(8.9),
        Inches(5.55),
    )
    add_box(
        slide,
        Inches(1.1),
        Inches(6.35),
        Inches(7.8),
        Inches(0.4),
        "从事实候选到 A/B 取舍，再到状态更新与解释输出。",
        fill=RGBColor(255, 251, 235),
        line=COLORS["orange"],
        font_size=12.5,
    )
    add_source(slide, "fact_ledger_v0.md §1-2")


def slide_s05(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "系统架构把 Data、Agent 和 Benchmark 串成可复查的推荐闭环",
        "数据证据层、智能体服务、推荐决策和评测环境各自承担边界",
        "05",
    )
    arch = generated_asset("architecture_wide.png")
    if arch and USE_ARCHITECTURE_IMAGE_ON_MAIN:
        add_image_fit(
            slide, arch, Inches(0.55), Inches(1.12), Inches(8.9), Inches(5.55)
        )
        add_box(
            slide,
            Inches(1.1),
            Inches(6.3),
            Inches(7.8),
            Inches(0.34),
            "数据给事实边界，Agent 组织交互，Benchmark 检查机制有效性",
            fill=COLORS["white"],
            line=COLORS["line"],
            font_size=10.5,
        )
        add_source(
            slide,
            "generated_assets/architecture_wide.png; fig_4_1_system_architecture semantics",
        )
        return
    rows = [
        ("Data / evidence", "PostgreSQL、专业树、地域树、质量画像", COLORS["blue"]),
        ("Agent service", "语义归一、探测规划、表达组织", COLORS["green"]),
        ("Decision loop", "SAVF、UCB、Pareto A/B、BT posterior", COLORS["orange"]),
        ("Benchmark / UI", "受控画像、自动对弈、可回放界面", COLORS["zju_blue"]),
    ]
    for idx, (head, body, color) in enumerate(rows):
        y = Inches(1.25 + idx * 1.18)
        add_box(
            slide,
            Inches(0.75),
            y,
            Inches(2.05),
            Inches(0.72),
            head,
            fill=COLORS["white"],
            line=color,
            font_size=14,
            bold=True,
        )
        add_box(
            slide,
            Inches(3.0),
            y,
            Inches(5.85),
            Inches(0.72),
            body,
            fill=RGBColor(248, 250, 252),
            line=COLORS["line"],
            font_size=13,
        )
        if idx < len(rows) - 1:
            add_arrow(
                slide,
                Inches(1.63),
                y + Inches(0.82),
                w=Inches(0.42),
                h=Inches(0.18),
                color=COLORS["line"],
            )
    add_box(
        slide,
        Inches(1.15),
        Inches(6.0),
        Inches(7.7),
        Inches(0.44),
        "关键边界：LLM 组织表达和探测，不直接生成事实候选",
        fill=RGBColor(255, 251, 235),
        line=COLORS["orange"],
        font_size=14,
        bold=True,
    )
    add_source(
        slide,
        "source-derived redraw from fig_4_1_system_architecture / render_system_architecture()",
    )


def slide_s06(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "事实候选由结构化证据层给出，LLM 不生成学校、专业、分数或学费",
        "事实边界先于推荐解释",
        "06",
    )
    add_pipeline(
        slide,
        ["用户表达", "LLM 语义归一", "PostgreSQL / 证据层", "事实候选集", "推荐解释"],
        Inches(2.15),
        x0=Inches(0.7),
        box_w=Inches(1.55),
        box_h=Inches(0.72),
        colors=[
            COLORS["line"],
            COLORS["green"],
            COLORS["blue"],
            COLORS["blue"],
            COLORS["zju_blue"],
        ],
        font_size=12,
    )
    add_box(
        slide,
        Inches(3.65),
        Inches(3.5),
        Inches(2.8),
        Inches(1.05),
        "只在确定性数据源中取学校、专业、分数、位次、学费",
        fill=RGBColor(241, 248, 255),
        line=COLORS["blue"],
        font_size=14,
        bold=True,
    )
    add_box(
        slide,
        Inches(1.0),
        Inches(5.2),
        Inches(8.0),
        Inches(0.5),
        "答辩安全句：事实由数据给出，表达由模型完成，取舍由用户确认",
        fill=COLORS["white"],
        line=COLORS["green"],
        font_size=15,
    )
    add_source(slide, "fact_ledger_v0.md §2, §5, §8")


def slide_s07(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "专业树和地域树提供可审计覆盖，但不被夸大为全部语义人工正确",
        "覆盖贡献和边界同时讲清",
        "07",
    )
    headers = ["证据层", "覆盖事实", "答辩边界"]
    rows = [
        [
            "专业层级本体",
            "22,759 / 22,759 专业名\n140,995 / 140,995 录取记录",
            "可审计挂载，不等于全部语义人工正确",
        ],
        [
            "地域层级画像",
            "414 省-市对\n3,219 所学校\n35 个省份映射",
            "地理层级证据，不编码就业或生活质量",
        ],
        [
            "验证集",
            "Accuracy 0.7169\nMacro-F1 0.7136\nHit@10 0.9819",
            "验证集结果不可外推为全库正确率",
        ],
    ]
    add_table_like(
        slide,
        Inches(0.55),
        Inches(1.25),
        Inches(8.9),
        Inches(0.86),
        headers,
        rows,
        col_ratios=[0.22, 0.34, 0.44],
        font_size=11.5,
    )
    add_source(slide, "fact_ledger_v0.md §4; major/region reports")


def slide_s08(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "静态检索能找到事实，却难以判断用户愿意为哪类收益放宽哪条约束",
        "静态基线解决事实召回，多轮闭环解决隐藏底线",
        "08",
    )
    add_image_fit(
        slide,
        cropped_figure("fig_3_1_v1_hybrid_rag_flow"),
        Inches(0.55),
        Inches(1.12),
        Inches(7.35),
        Inches(5.45),
    )
    add_box(
        slide,
        Inches(8.05),
        Inches(1.35),
        Inches(1.35),
        Inches(1.25),
        "强项\n事实召回\n候选排序",
        fill=RGBColor(241, 248, 255),
        line=COLORS["blue"],
        font_size=12.5,
        bold=True,
    )
    add_box(
        slide,
        Inches(8.05),
        Inches(3.0),
        Inches(1.35),
        Inches(1.25),
        "缺口\n不主动问\n隐藏底线",
        fill=RGBColor(255, 248, 240),
        line=COLORS["orange"],
        font_size=12.5,
        bold=True,
    )
    add_box(
        slide,
        Inches(8.05),
        Inches(4.65),
        Inches(1.35),
        Inches(1.25),
        "推进\nA/B取舍\n后验状态",
        fill=RGBColor(241, 250, 246),
        line=COLORS["green"],
        font_size=12.5,
        bold=True,
    )
    add_source(slide, "fact_ledger_v0.md §5; 02-problem-algorithm.tex:81-103")


def slide_s09(prs):
    slide = add_slide(prs)
    generated = vai01_asset()
    if generated and USE_VAI01_ON_MAIN:
        add_title(
            slide,
            "推荐决策闭环把事实候选转化为可回答的偏好取舍",
            "SAVF、UCB、Pareto/BT 与后验状态组成可复查的多轮闭环",
            "09",
        )
        add_image_fit(
            slide, generated, Inches(0.55), Inches(1.08), Inches(8.9), Inches(5.65)
        )
        add_box(
            slide,
            Inches(1.05),
            Inches(6.35),
            Inches(7.9),
            Inches(0.34),
            "解释性示意图：图中候选、权重和示例均为抽象占位，不承载新增实验事实",
            fill=COLORS["white"],
            line=COLORS["line"],
            font_size=10.5,
        )
        add_source(
            slide,
            "generated_assets/"
            + vai01_asset_label()
            + "; academic_figure_prompt_v0.md",
        )
        return
    add_title(
        slide,
        "推荐决策闭环把事实候选转化为可回答的偏好取舍",
        "SAVF、UCB、Pareto/BT 与后验状态组成可复查的多轮闭环",
        "09",
    )
    labels = [
        "事实候选",
        "SAVF\n保护底线",
        "UCB\n选轴",
        "Pareto A/B\n具体取舍",
        "用户反馈",
        "BT 后验\n偏好状态",
    ]
    add_pipeline(
        slide,
        labels,
        Inches(1.58),
        x0=Inches(0.55),
        box_w=Inches(1.35),
        box_h=Inches(0.8),
        colors=[
            COLORS["blue"],
            COLORS["orange"],
            COLORS["green"],
            COLORS["orange"],
            COLORS["line"],
            COLORS["zju_blue"],
        ],
        font_size=10.5,
    )
    add_box(
        slide,
        Inches(1.0),
        Inches(3.25),
        Inches(2.15),
        Inches(1.1),
        "protect\n严重违约不能被高分抵消",
        fill=RGBColor(255, 248, 240),
        line=COLORS["orange"],
        font_size=13,
        bold=True,
    )
    add_box(
        slide,
        Inches(3.95),
        Inches(3.25),
        Inches(2.15),
        Inches(1.1),
        "choose\n优先问最有诊断价值的维度",
        fill=RGBColor(241, 250, 246),
        line=COLORS["green"],
        font_size=13,
        bold=True,
    )
    add_box(
        slide,
        Inches(6.9),
        Inches(3.25),
        Inches(2.15),
        Inches(1.1),
        "remember\n反馈写回在线偏好状态",
        fill=RGBColor(241, 248, 255),
        line=COLORS["blue"],
        font_size=13,
        bold=True,
    )
    add_box(
        slide,
        Inches(1.2),
        Inches(5.45),
        Inches(7.6),
        Inches(0.52),
        "这张图是解释性视觉，不承载新增实验事实",
        fill=COLORS["white"],
        line=COLORS["line"],
        font_size=13,
    )
    add_source(slide, "fact_ledger_v0.md §5")


def slide_s10(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "非补偿性价值映射防止预算、专业等底线被线性总分掩盖",
        "严重违约先被惩罚，避免被学校层次等高分抵消",
        "10",
    )
    add_image_fit(
        slide,
        generated_asset("savf_mechanism_ppt.png"),
        Inches(0.55),
        Inches(1.1),
        Inches(8.9),
        Inches(5.65),
    )
    add_box(
        slide,
        Inches(2.55),
        Inches(5.78),
        Inches(4.95),
        Inches(0.44),
        "属性值 -> 违约惩罚 -> 底线保护后的评分",
        fill=COLORS["white"],
        line=COLORS["blue"],
        font_size=13,
        bold=True,
        radius=False,
    )
    add_source(slide, "fact_ledger_v0.md §5; 02-problem-algorithm.tex:119-152")


def slide_s11(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "UCB 主动探测让系统优先询问最有诊断价值的偏好维度",
        "收益线索和不确定性共同决定下一轮探测维度",
        "11",
    )
    img = cropped_figure("fig_5_3_ucb_dispatch")
    add_image_fit(slide, img, Inches(0.55), Inches(1.08), Inches(6.25), Inches(5.75))
    add_box(
        slide,
        Inches(7.05),
        Inches(1.45),
        Inches(2.25),
        Inches(0.75),
        "收益线索",
        fill=COLORS["white"],
        line=COLORS["blue"],
        font_size=15,
        bold=True,
    )
    add_box(
        slide,
        Inches(7.05),
        Inches(2.55),
        Inches(2.25),
        Inches(0.75),
        "当前不确定性",
        fill=COLORS["white"],
        line=COLORS["green"],
        font_size=15,
        bold=True,
    )
    add_arrow(
        slide,
        Inches(7.95),
        Inches(3.55),
        w=Inches(0.42),
        h=Inches(0.22),
        color=COLORS["line"],
    )
    add_box(
        slide,
        Inches(7.05),
        Inches(4.05),
        Inches(2.25),
        Inches(0.85),
        "下一轮探测维度",
        fill=RGBColor(241, 250, 246),
        line=COLORS["green"],
        font_size=15,
        bold=True,
    )
    add_source(slide, "fig_5_3_ucb_dispatch; fact_ledger_v0.md §5")


def slide_s12(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "帕累托 A/B 候选把抽象偏好变成用户能够回答的具体取舍",
        "BT 后验写回偏好状态，而不是直接改写 SQL 参数",
        "12",
    )
    add_box(
        slide,
        Inches(0.75),
        Inches(1.35),
        Inches(3.0),
        Inches(2.0),
        "候选 A\n风险更稳\n专业匹配较弱\n学费更低",
        fill=RGBColor(241, 248, 255),
        line=COLORS["blue"],
        font_size=15,
        bold=True,
    )
    add_box(
        slide,
        Inches(4.0),
        Inches(1.35),
        Inches(3.0),
        Inches(2.0),
        "候选 B\n专业更匹配\n风险略高\n可能需要放宽预算",
        fill=RGBColor(255, 248, 240),
        line=COLORS["orange"],
        font_size=15,
        bold=True,
    )
    add_box(
        slide,
        Inches(7.35),
        Inches(1.35),
        Inches(1.85),
        Inches(2.0),
        "反馈\nPrefer A\nPrefer B\nUncertain",
        fill=COLORS["white"],
        line=COLORS["line"],
        font_size=13,
    )
    add_box(
        slide,
        Inches(2.0),
        Inches(4.25),
        Inches(6.0),
        Inches(0.9),
        "BT posterior: region | major | tuition | risk",
        fill=RGBColor(248, 250, 252),
        line=COLORS["zju_blue"],
        font_size=15,
        bold=True,
    )
    for i, color in enumerate(
        [COLORS["blue"], COLORS["green"], COLORS["orange"], COLORS["line"]]
    ):
        add_box(
            slide,
            Inches(2.35 + i * 1.33),
            Inches(5.35),
            Inches(0.9),
            Inches(0.34),
            "",
            fill=color,
            line=color,
            font_size=8,
            radius=False,
        )
    add_source(slide, "fact_ledger_v0.md §5; 03-system-design.tex:317-339")


def slide_s13(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "运行时状态机把澄清、反馈和终局推荐串成可回放链路",
        "每轮问答、反馈和状态更新都有明确阶段",
        "13",
    )
    img = cropped_figure("fig_5_2_runtime_state_machine")
    add_image_fit(slide, img, Inches(0.55), Inches(1.08), Inches(6.9), Inches(5.75))
    add_bullets(
        slide,
        Inches(7.65),
        Inches(1.35),
        Inches(1.8),
        Inches(3.9),
        [
            "探索期和终局推荐期分离",
            "interrupt / resume 支持人机交互",
            "每轮为什么问、怎么答、状态如何变都可回放",
        ],
        size=11.5,
    )
    add_source(slide, "fig_5_2_runtime_state_machine; fact_ledger_v0.md §3, §5")


def slide_s14(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "前端把偏好澄清过程呈现为可理解的取舍问题和推荐报告",
        "用户看到的是具体取舍和可追溯解释",
        "14",
    )
    add_image_fit(
        slide,
        image_path("fig_3_5_elicitation_console"),
        Inches(0.55),
        Inches(1.22),
        Inches(4.25),
        Inches(4.75),
    )
    add_image_fit(
        slide,
        image_path("fig_3_6_final_decision_report"),
        Inches(5.2),
        Inches(1.22),
        Inches(4.25),
        Inches(4.75),
    )
    add_box(
        slide,
        Inches(0.9),
        Inches(6.15),
        Inches(3.55),
        Inches(0.36),
        "取舍问题",
        fill=COLORS["white"],
        line=COLORS["blue"],
        font_size=12,
    )
    add_box(
        slide,
        Inches(5.65),
        Inches(6.15),
        Inches(3.35),
        Inches(0.36),
        "推荐解释",
        fill=COLORS["white"],
        line=COLORS["green"],
        font_size=12,
    )
    add_source(slide, "fig_3_5_elicitation_console; fig_3_6_final_decision_report")


def slide_s15(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "受控测试集用隐藏底线用户画像检验系统是否真的恢复真实偏好",
        "隐藏画像只供评测端使用，不进入 Agent 输入",
        "15",
    )
    add_image_fit(
        slide,
        cropped_figure("fig_4_2_benchmark_flow"),
        Inches(0.55),
        Inches(1.08),
        Inches(6.85),
        Inches(5.75),
    )
    add_box(
        slide,
        Inches(7.65),
        Inches(1.45),
        Inches(1.75),
        Inches(0.78),
        "180 条\n受控画像",
        fill=RGBColor(241, 248, 255),
        line=COLORS["blue"],
        font_size=16,
        bold=True,
    )
    add_box(
        slide,
        Inches(7.65),
        Inches(2.65),
        Inches(1.75),
        Inches(0.78),
        "显式话语\n+ 隐藏底线",
        fill=RGBColor(255, 248, 240),
        line=COLORS["orange"],
        font_size=16,
        bold=True,
    )
    add_box(
        slide,
        Inches(7.65),
        Inches(3.85),
        Inches(1.75),
        Inches(0.78),
        "F1@N / MAE\n过程指标",
        fill=RGBColor(241, 250, 246),
        line=COLORS["green"],
        font_size=16,
        bold=True,
    )
    add_source(slide, "fig_4_2_benchmark_flow; fact_ledger_v0.md §6.1")


def slide_result(prs, num, title, img_name, callout, source):
    slide = add_slide(prs)
    add_title(slide, title, "同一评价口径下比较完整系统、静态检索提示和消融系统", num)
    add_image_fit(
        slide,
        image_path(img_name),
        Inches(0.65),
        Inches(1.05),
        Inches(8.7),
        Inches(5.15),
    )
    add_takeaway(slide, callout.replace("\n", "；"), y=Inches(6.28))
    add_source(slide, source)


def slide_s18(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "过程指标说明完整系统更能形成有张力的问题并稳定更新偏好状态",
        "过程证据解释主动探测和后验追踪的作用",
        "18",
    )
    add_image_fit(
        slide,
        image_path("fig_4_8_1_c1_planner_process"),
        Inches(0.55),
        Inches(1.08),
        Inches(7.15),
        Inches(5.75),
    )
    add_pipeline(
        slide,
        ["选轴", "取舍", "后验"],
        Inches(2.1),
        x0=Inches(7.95),
        box_w=Inches(0.72),
        box_h=Inches(0.52),
        gap=Inches(0.08),
        colors=[COLORS["green"], COLORS["orange"], COLORS["blue"]],
        font_size=10,
    )
    add_box(
        slide,
        Inches(7.85),
        Inches(3.55),
        Inches(1.65),
        Inches(1.3),
        "过程指标解释为什么主动探测与后验追踪有效",
        fill=COLORS["white"],
        line=COLORS["line"],
        font_size=11.5,
    )
    add_source(slide, "fig_4_8_1_c1_planner_process; fact_ledger_v0.md §6.1")


def slide_s19(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "本设计把事实证据、主动澄清、在线后验和可复查评测组织成一个工程闭环",
        "事实、交互、状态和评测共同支撑可复查推荐",
        "19",
    )
    labels = [
        "问题定义\n隐藏底线",
        "数据证据\n可审计候选",
        "推荐决策\n主动取舍",
        "运行链路\n可回放",
        "受控测试\n机制验证",
    ]
    add_pipeline(
        slide,
        labels,
        Inches(2.15),
        x0=Inches(0.6),
        box_w=Inches(1.55),
        box_h=Inches(0.86),
        colors=[
            COLORS["zju_blue"],
            COLORS["blue"],
            COLORS["orange"],
            COLORS["green"],
            COLORS["zju_blue"],
        ],
        font_size=11.5,
    )
    add_box(
        slide,
        Inches(1.25),
        Inches(4.25),
        Inches(7.5),
        Inches(0.82),
        "回扣开场：高风险推荐要同时守住事实可信和偏好真实",
        fill=RGBColor(248, 250, 252),
        line=COLORS["line"],
        font_size=17,
        bold=True,
    )
    add_source(slide, "fact_ledger_v0.md §3; defense_narrative_v0.md §4.5")


def slide_s20(prs):
    slide = add_slide(prs)
    add_title(
        slide,
        "系统仍需真实用户研究和泛化验证，但已初步支持先问清底线再推荐的技术路径",
        "当前结论是机制验证，不是生产级有效性证明",
        "20",
    )
    headers = ["当前边界", "后续工作"]
    rows = [
        ["受控测试不能替代真实学生或专家评估", "小规模真实用户/专家评估"],
        ["实验范围仍需跨省、跨年份泛化", "扩展年份、省份和证据维度"],
        ["用户控制与隐私合规仍需强化", "更完善的控制面板与合规流程"],
    ]
    add_table_like(
        slide,
        Inches(0.85),
        Inches(1.35),
        Inches(8.3),
        Inches(0.82),
        headers,
        rows,
        col_ratios=[0.5, 0.5],
        font_size=13,
    )
    add_box(
        slide,
        Inches(1.25),
        Inches(5.5),
        Inches(7.5),
        Inches(0.5),
        "当前贡献：把事实、取舍、反馈和后验状态连成可复查机制",
        fill=RGBColor(241, 250, 246),
        line=COLORS["green"],
        font_size=14,
        bold=True,
    )
    add_source(slide, "fact_ledger_v0.md §1, §8; 07-conclusion.tex:28-38")


def slide_s21(prs):
    slide = add_slide(prs)
    add_box(
        slide,
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(7.5),
        "",
        fill=COLORS["white"],
        line=COLORS["white"],
        radius=False,
    )
    add_text(
        slide,
        Inches(0.9),
        Inches(2.2),
        Inches(8.2),
        Inches(0.7),
        "谢谢各位老师",
        size=34,
        bold=True,
        color=COLORS["zju_blue"],
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        Inches(0.9),
        Inches(3.35),
        Inches(8.2),
        Inches(0.45),
        "欢迎批评指正",
        size=24,
        color=COLORS["zju_red"],
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        Inches(0.9),
        Inches(5.7),
        Inches(8.2),
        Inches(0.3),
        "Q&A",
        size=18,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
    )


def add_backup_separator(prs):
    slide = add_slide(prs)
    add_box(
        slide,
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(7.5),
        "",
        fill=COLORS["white"],
        line=COLORS["white"],
        radius=False,
    )
    add_text(
        slide,
        Inches(0.8),
        Inches(2.35),
        Inches(8.4),
        Inches(0.68),
        "Backup Slides",
        size=32,
        bold=True,
        color=COLORS["zju_blue"],
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        Inches(0.8),
        Inches(3.25),
        Inches(8.4),
        Inches(0.4),
        "答辩问答备份材料",
        size=20,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
    )
    return slide


def backup_slide(prs, sid, title, bullets, img_name=None, table=None):
    slide = add_slide(prs)
    add_title(slide, title, "Backup " + sid, sid)
    if img_name:
        add_image_fit(
            slide,
            image_path(img_name),
            Inches(0.65),
            Inches(1.22),
            Inches(5.25),
            Inches(5.35),
        )
        add_bullets(
            slide,
            Inches(6.15),
            Inches(1.3),
            Inches(3.1),
            Inches(4.9),
            bullets,
            size=12.5,
        )
    elif table:
        add_table_like(
            slide,
            Inches(0.65),
            Inches(1.35),
            Inches(8.7),
            Inches(0.72),
            table["headers"],
            table["rows"],
            col_ratios=table.get("ratios"),
            font_size=10.5,
        )
    else:
        add_bullets(
            slide,
            Inches(0.95),
            Inches(1.35),
            Inches(8.0),
            Inches(4.85),
            bullets,
            size=14,
        )
    add_source(slide, "ppt_defense_qa_backup_v0.md; fact_ledger_v0.md")


def backup_b16_algorithm_visual(prs):
    slide = add_slide(prs)
    add_title(slide, "算法示意图只解释机制边界，不作为新增证据", "Backup B16", "B16")
    bullets = [
        "总览 SAVF、UCB、Pareto/BT 与后验状态如何形成闭环。",
        "图中候选、权重和示例均为抽象占位。",
        "不含真实学校名、分数、产品 UI 或实验数值。",
    ]
    add_bullets(
        slide, Inches(0.7), Inches(1.35), Inches(3.0), Inches(2.45), bullets, size=13
    )
    generated = vai01_asset()
    if generated:
        add_image_fit(
            slide, generated, Inches(3.85), Inches(1.25), Inches(5.75), Inches(4.35)
        )
    add_box(
        slide,
        Inches(0.85),
        Inches(5.45),
        Inches(8.3),
        Inches(0.42),
        "答辩口径：它帮助解释机制，不替代论文实验图和结果表。",
        fill=RGBColor(255, 251, 235),
        line=COLORS["orange"],
        font_size=12.5,
    )


def add_backups(prs):
    add_backup_separator(prs)
    backup_slide(
        prs,
        "B01",
        "LLM 不生成事实候选，事实边界由证据层控制",
        [
            "LLM 负责语义归一、探测规划和表达组织。",
            "学校、专业、分数、位次、学费来自 PostgreSQL 和标准化证据层。",
            "安全回答：事实由数据给出，取舍由用户确认。",
        ],
    )
    backup_slide(
        prs,
        "B02",
        "专业树覆盖是可审计挂载，不等于全语义正确",
        [],
        table={
            "headers": ["指标", "数值", "边界"],
            "rows": [
                ["原始专业名", "22,759 / 22,759", "覆盖不等于全语义正确"],
                ["录取记录", "140,995 / 140,995", "可审计挂载"],
                [
                    "Clean validation",
                    "Accuracy 0.7169 / Macro-F1 0.7136",
                    "不可外推为全库正确率",
                ],
                ["剩余错分", "47 / 166", "后续 HITL 优化空间"],
            ],
            "ratios": [0.28, 0.34, 0.38],
        },
    )
    backup_slide(
        prs,
        "B03",
        "地域树提供地理层级证据，不编码就业或生活质量",
        [
            "414 个省-市对。",
            "3,219 所学校。",
            "35 个省份映射，remaining unassigned = 0。",
            "不能把地域层级直接等价为就业机会、生活成本或城市生活质量。",
        ],
        img_name="fig_4_5_region_hierarchy_partial",
    )
    backup_slide(
        prs,
        "B04",
        "静态检索是强基线，但不能自然恢复隐藏底线",
        [
            "显式约束归一、数据库过滤、向量召回和重排构成强工程基线。",
            "缺口不是事实召回，而是用户愿意为哪类收益放宽哪条底线。",
            "不贬低 RAG，只说明本文推进点。",
        ],
        img_name="fig_3_1_v1_hybrid_rag_flow",
    )
    backup_slide(
        prs,
        "B05",
        "SAVF 保护底线，防止线性总分掩盖严重违约",
        [
            "预算、专业、地域等底线先做单属性价值映射。",
            "严重违约不应被学校层次等高分维度线性抵消。",
            "公式只作辅助，主讲解释非补偿性直觉。",
        ],
    )
    backup_slide(
        prs,
        "B06",
        "UCB 是工程启发式选轴，不是严格最优证明",
        [
            "结合收益线索和当前不确定性。",
            "目标是在有限轮次优先问更可能暴露隐藏底线的维度。",
            "后续可比较随机选轴、最大方差、最大均值等策略。",
        ],
        img_name="fig_5_3_ucb_dispatch",
    )
    backup_slide(
        prs,
        "B07",
        "Pareto A/B 与 BT 后验把取舍反馈写入偏好状态",
        [
            "A/B 候选展示边际替代关系。",
            "用户只需表达接受、拒绝或偏向哪一边。",
            "BT 后验影响后续选轴、候选评分或解释，不直接改写 SQL。",
        ],
    )
    backup_slide(
        prs,
        "B08",
        "运行时状态机让澄清和终局推荐可回放",
        [
            "隔离探索期和终局推荐期。",
            "支持 interrupt / resume。",
            "每轮提问、反馈和状态变化可以审计。",
        ],
        img_name="fig_5_2_runtime_state_machine",
    )
    backup_slide(
        prs,
        "B09",
        "Benchmark 信息边界：Agent 不读取 hidden persona",
        [
            "Simulator 和 evaluator 使用隐藏字段。",
            "Agent 只看到显式话语与可检索证据。",
            "这样才能检查系统是否在多轮过程中恢复隐藏底线。",
        ],
    )
    backup_slide(
        prs,
        "B10",
        "基线横评和主消融用不同问题回答机制贡献",
        [
            "横评：五个基座模型，比较完整系统与静态检索提示。",
            "消融：固定 GLM-5.1，去除主动探测或后验追踪。",
            "回答的是链路与模块贡献，不是生产级泛化。",
        ],
        img_name="fig_4_5_c1_baseline_model_target",
    )
    backup_slide(
        prs,
        "B11",
        "过程指标解释“选轴 -> 取舍 -> 后验”为什么起作用",
        [
            "有效探测方向。",
            "权衡张力。",
            "后验状态更新。",
            "过程图只补充机制解释，主线可跳过。",
        ],
        img_name="fig_4_8_2_c1_negotiator_process",
    )
    backup_slide(
        prs,
        "B12",
        "旧口径实验只作项目演化备份，不覆盖终稿主线",
        [
            "当前答辩以最终 LaTeX 第 4 章 180 条受控画像实验为准。",
            "major_geo_v1、risk_band_v1 等旧口径只解释项目演化。",
            "multi_axis_v2 是压力测试，不是第八组主实验。",
        ],
    )
    backup_slide(
        prs,
        "B13",
        "局限与后续工作：真实用户、跨省跨年份、用户控制",
        [
            "当前是机制验证，不替代真实用户研究。",
            "后续需要专家评估、跨省跨年份泛化和隐私合规。",
            "先说明当前机制价值，再说明边界。",
        ],
    )
    backup_slide(
        prs,
        "B14",
        "现有产品简要现状评测不等于竞品报告",
        [
            "只用作 20-30 秒场景入口。",
            "若使用截图，必须记录 URL 和访问日期。",
            "不把阳光高考、夸克高考作为实验基线。",
        ],
    )
    backup_slide(
        prs,
        "B15",
        "工程实现路径：四层架构与关键模块归属",
        [
            "Data 提供事实边界。",
            "Agent 组织语义归一和探测规划。",
            "Decision loop 完成取舍与偏好状态更新。",
            "Benchmark 检验机制有效性。",
        ],
    )
    if generated_asset("architecture_wide.png"):
        slide = prs.slides[-1]
        add_image_fit(
            slide,
            generated_asset("architecture_wide.png"),
            Inches(0.65),
            Inches(4.05),
            Inches(8.7),
            Inches(2.55),
        )
    backup_b16_algorithm_visual(prs)


def add_main_slides(prs):
    make_cover(prs)
    add_agenda_slide(prs)
    add_section_divider(prs, "一", "研究背景与问题定义")
    slide_problem_combined(prs)
    slide_s04(prs)
    add_section_divider(prs, "二", "系统与方法设计")
    slide_s05(prs)
    slide_fact_boundary_compact(prs)
    slide_s08(prs)
    slide_s09(prs)
    slide_s10(prs)
    slide_s11(prs)
    slide_s12_state_combined(prs)
    add_section_divider(prs, "三", "实验验证与性能分析")
    slide_s15(prs)
    slide_result(
        prs,
        "16",
        "参考基线横评显示，完整系统的优势来自交互式证据链路而不是更强的单轮生成",
        "fig_4_5_c1_baseline_model_target",
        "整体优于静态提示\n但只表述为趋势支持",
        "fig_4_5_c1_baseline_model_target; fact_ledger_v0.md §6.1",
    )
    slide_result(
        prs,
        "17",
        "主消融实验显示，主动探测和后验追踪分别支撑探测方向选择与反馈吸收",
        "fig_4_6_c1_ablation_core_metrics",
        "去除模块后退化\n机制贡献得到初步支持",
        "fig_4_6_c1_ablation_core_metrics; fact_ledger_v0.md §6.1",
    )
    slide_conclusion_combined(prs)
    slide_s21(prs)


def write_manifest(slide_count, used_images):
    rel_images = []
    for item in used_images:
        try:
            rel_images.append(str(Path(item).relative_to(ROOT)).replace("\\", "/"))
        except ValueError:
            rel_images.append(str(item))
    manifest = """---
stage: deck_build
stage_status: draft
requires_confirmed:
  - ppt_production_brief
  - fact_ledger
  - defense_narrative
  - storyboard
  - speaker_notes_rehearsal
  - defense_qa_backup
  - asset_layout_plan
  - academic_figure_prompt_when_required
  - content_fidelity_qa
allowed_next_stage: ppt-render-qa-loop
confirmed_by:
created_at: 2026-05-31
---

# PPT deck build manifest v0

## 1. Output

- PPTX path: `{deck_path}`
- Build script: `exp/ppt_deck_build_v0/build_deck_v0.py`
- Local figure generation script: `exp/ppt_deck_build_v0/generate_local_figures_v0.py`
- Cropped source-figure preparation script: `exp/ppt_deck_build_v0/prepare_cropped_assets_v0.py`
- Stage status: draft
- Main slides: 19
- Backup separator: 1
- Backup slides: 16
- Total slides: {slide_count}
- Aspect ratio: 4:3, 10.0 x 7.5 in

## 2. Confirmed input artifacts

- `align/ppt_production_brief_v0.md`: confirmed
- `align/fact_ledger_v0.md`: confirmed
- `align/ppt_defense_narrative_v0.md`: confirmed
- `align/PPT_storyboard_v0.md`: confirmed
- `align/ppt_speaker_notes_rehearsal_v0.md`: confirmed
- `align/ppt_defense_qa_backup_v0.md`: confirmed
- `align/PPT_asset_audit_v0.md`: confirmed
- `align/template_inventory_v0.md`: confirmed
- `align/template_design_rules_v0.md`: confirmed
- `align/ppt_layout_plan_v0.json`: confirmed
- `align/academic_figure_prompt_v0.md`: confirmed
- `align/ppt_content_fidelity_qa_v0.md`: confirmed
- `align/ppt_deck_visual_refactor_plan_v0.md`: confirmed

## 3. Visual route outcomes

| Route | Outcome |
| --- | --- |
| Template | Used the confirmed 4:3 slide size and template policy. The generation script starts from `zjuslides.pptx` when available, then builds custom editable slides. |
| Main evidence figures | Inserted source PNG figures from the final thesis/project figure directories where appropriate. |
| Figure scaling | Created non-destructive cropped copies of source diagrams under `generated_assets/crop_*.png` to remove large white margins and make main-slide figures larger. |
| Local mechanism figure | Added `generated_assets/savf_mechanism_ppt.png` for Slide 10 to replace scattered text boxes with a single mechanism visual. |
| Visual refactor | Applied `align/ppt_deck_visual_refactor_plan_v0.md`: 19 main slides, agenda/section dividers, light title band, and one-main-visual evidence grammar. |
| Editable diagrams | Built the problem framing, fact-boundary, method, A/B posterior, conclusion, and several backup slides with editable PowerPoint shapes/tables. |
| V-AI01 generated academic figure | `generated_assets/v_ai01_openrouter_icu.png` is used on Slide 9 as the main method overview and repeated in B16 as an explanation boundary backup. If it is missing, the local deterministic fallback `generated_assets/v_ai01_algorithm_flow.png` is used. |
| Product screenshots | Not used. Slide 2 and B14 use abstract comparison to avoid stale product claims. |
| Horizontal architecture | Main Slide 5 keeps an editable source-derived architecture summary. If `generated_assets/architecture_wide.png` exists, it is added to B15 by default; the thesis vertical figure is not overwritten or force-fit. |

User feedback applied: V-AI01 is now visible in the main talk on Slide 9, with B16 retained as the backup/boundary explanation slide. Internal artifact/source footer lines are hidden from audience-facing slides.

OpenRouter ICU status: completed for `generated_assets/v_ai01_openrouter_icu.png` when that file is present. It was generated from the confirmed `align/academic_figure_prompt_v0.md` prompt with `gpt-image-2`, size `1536x1152`, quality `medium`, output `png`; request metadata is recorded in `exp/ppt_deck_build_v0/openrouter_vai01_generation_log.json`. Partial images, if present, are intermediate artifacts and are not inserted into the deck.

## 4. Notes and editability

- Speaker notes were not inserted into the PowerPoint notes pane because `python-pptx` does not support notes-pane authoring reliably in this environment.
- Confirmed notes remain in `align/ppt_speaker_notes_rehearsal_v0.md`.
- Most text, tables, callouts, and diagrams are editable PowerPoint objects.
- Source evidence figures and UI screenshots are raster PNG insertions with editable surrounding labels/callouts.

## 5. Assets used

{assets}

## 6. Known pre-render risks

- Visual polish has not been render-QA checked yet.
- Some source figure labels may be too small at 4:3 slideshow size, especially UI screenshots and dense experiment charts.
- Backup slides are placed after a Backup separator rather than hidden, so render QA should verify navigation/visibility policy.
- Slide 5 is a compact PPT redraw of the architecture. Render QA should compare it against the original semantics.
- The previous process-indicator Slide 18 was moved out of the main talk path; detailed process evidence remains available in backup.
- Notes-pane insertion is not available; presenter should use the confirmed notes artifact.

## 7. Handoff

If the user accepts this draft for visual QA, update only this manifest to:

```yaml
stage_status: confirmed
confirmed_by: user, 2026-05-31
```

Then stop. The next stage is `ppt-render-qa-loop`.
""".format(
        deck_path=str(DECK_PATH).replace("\\", "/"),
        slide_count=slide_count,
        assets="\n".join("- `" + p + "`" for p in rel_images)
        if rel_images
        else "- No raster assets were inserted.",
    )
    MANIFEST_PATH.write_text(manifest, encoding="utf-8")


def collect_used_images(prs):
    # Re-scan relationships after save would be expensive; record all source images that exist.
    names = [
        "fig_5_3_ucb_dispatch",
        "fig_5_2_runtime_state_machine",
        "fig_3_5_elicitation_console",
        "fig_3_6_final_decision_report",
        "fig_4_2_benchmark_flow",
        "fig_4_5_c1_baseline_model_target",
        "fig_4_6_c1_ablation_core_metrics",
        "fig_4_8_1_c1_planner_process",
        "fig_4_5_region_hierarchy_partial",
        "fig_3_1_v1_hybrid_rag_flow",
        "fig_4_8_2_c1_negotiator_process",
    ]
    paths = []
    for name in names:
        path = image_path(name)
        if path:
            paths.append(path)
    generated_names = [
        "crop_fig_5_1_mas_workflow.png",
        "crop_fig_3_1_v1_hybrid_rag_flow.png",
        "crop_fig_5_3_ucb_dispatch.png",
        "crop_fig_5_2_runtime_state_machine.png",
        "crop_fig_4_2_benchmark_flow.png",
        "savf_mechanism_ppt.png",
        "architecture_wide.png",
    ]
    for path in [vai01_asset()] + [generated_asset(name) for name in generated_names]:
        if path:
            paths.append(path)
    return paths


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    if TEMPLATE.exists():
        prs = Presentation(str(TEMPLATE))
        remove_all_slides(prs)
    else:
        prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    add_main_slides(prs)
    add_backups(prs)
    used_images = collect_used_images(prs)
    prs.save(str(DECK_PATH))
    write_manifest(len(prs.slides), used_images)
    print("saved", DECK_PATH)
    print("slides", len(prs.slides))
    print("manifest", MANIFEST_PATH)


if __name__ == "__main__":
    main()
